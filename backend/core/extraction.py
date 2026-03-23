import os
import io

def extract_text(file_path: str, original_filename: str) -> str:
    """
    Extracts text from a given file path based on its extension.
    Supported types: .pdf, .pptx, .docx, .txt, .csv
    """
    ext = os.path.splitext(original_filename)[1].lower()
    
    if ext == ".txt":
        with open(file_path, "r", encoding="utf-8") as f:
            return f.read()
            
    elif ext == ".pdf":
        import pdfplumber
        text_content = []
        with pdfplumber.open(file_path) as pdf:
            for page in pdf.pages:
                page_text = page.extract_text()
                if page_text:
                    text_content.append(page_text)
        return "\n".join(text_content)
            
    elif ext == ".pptx":
        from pptx import Presentation
        prs = Presentation(file_path)
        text_content = []
        for slide in prs.slides:
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    text_content.append(shape.text)
        return "\n".join(text_content)

    elif ext == ".docx":
        import docx
        doc = docx.Document(file_path)
        return "\n".join([para.text for para in doc.paragraphs])

    elif ext == ".csv":
        import csv
        with open(file_path, "r", encoding="utf-8") as f:
            reader = csv.reader(f)
            return "\n".join([",".join(row) for row in reader])
            
    raise ValueError(f"Unsupported file extension: {ext}")

def chunk_text(text: str, chunk_size: int = 500, overlap: int = 50) -> list[str]:
    """
    Splits text into overlapping chunks, attempting to respect word boundaries.
    """
    if not text:
        return []
        
    words = text.split()
    chunks = []
    current_chunk_words = []
    current_length = 0
    
    # Simple word-based chunker with approximate character count
    for word in words:
        current_chunk_words.append(word)
        current_length += len(word) + 1 # +1 for space
        
        if current_length >= chunk_size:
            chunks.append(" ".join(current_chunk_words))
            # Overlap: keep some words for the next chunk
            overlap_words = []
            overlap_len = 0
            for w in reversed(current_chunk_words):
                if overlap_len + len(w) + 1 > overlap:
                    break
                overlap_words.insert(0, w)
                overlap_len += len(w) + 1
            
            current_chunk_words = overlap_words
            current_length = overlap_len
            
    if current_chunk_words:
        chunks.append(" ".join(current_chunk_words))
        
    return [c for c in chunks if c.strip()]

